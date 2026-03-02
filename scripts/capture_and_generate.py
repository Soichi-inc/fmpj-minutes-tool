#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Capture screenshots from the running dev server and generate PPTX manual."""

import os
import time
from playwright.sync_api import sync_playwright

ASSETS_DIR = "/Users/akirahasama/Desktop/dev/fmpj-minutes-tool/manual-assets"
BASE_URL = "http://127.0.0.1:3457"
PASSWORD = "fmpj2026"

os.makedirs(ASSETS_DIR, exist_ok=True)


def capture_screenshots():
    """Capture all necessary screenshots."""
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(
            viewport={"width": 1280, "height": 800},
            device_scale_factor=2,
        )
        page = context.new_page()

        # 1. Login page
        print("📸 Capturing login page...")
        page.goto(BASE_URL)
        page.wait_for_load_state("networkidle")
        time.sleep(1)
        page.screenshot(path=f"{ASSETS_DIR}/01_login.png")

        # 2. Login
        print("🔑 Logging in...")
        page.fill("#password", PASSWORD)
        page.click('button[type="submit"]')
        page.wait_for_url("**/dashboard**", timeout=10000)
        page.wait_for_load_state("networkidle")
        time.sleep(1)

        # 3. Step 1 - Empty dashboard
        print("📸 Capturing Step 1 (empty)...")
        page.screenshot(path=f"{ASSETS_DIR}/02_step1_empty.png")

        # 4. Fill Step 1 with sample data
        print("📝 Filling Step 1...")
        page.fill('input[placeholder*="第○回理事会"]', '第10回理事会')

        page.fill('input[type="date"]', '2026-03-01')
        page.fill('input[placeholder*="連盟会議室"]', '連盟会議室')

        # Fill attendees
        attendees_textarea = page.locator('textarea').first
        if attendees_textarea.count() > 0:
            attendees_textarea.fill('山田太郎\n佐藤花子\n田中一郎')

        # Fill transcript
        transcript_text = "Speaker 1 00:00\nそれでは第10回理事会を開催いたします。まず前回の議事録の承認から始めたいと思います。\n\nSpeaker 2 01:30\n前回の議事録について確認しましたが、特に修正事項はありません。\n\nSpeaker 3 02:15\n私も異議ありません。\n\nSpeaker 1 03:00\nありがとうございます。それでは前回議事録を承認とします。続いて来年度予算案について佐藤さんからお願いします。\n\nSpeaker 2 04:00\n来年度予算案について報告いたします。総額は前年度比5%増の1億2000万円を計画しています。主な増額要因として、デジタル化推進費用500万円、会員向けセミナー開催費用200万円の増額を見込んでいます。\n\nSpeaker 3 06:30\nデジタル化の具体的な内容を教えていただけますか。\n\nSpeaker 2 07:00\n会員管理システムのリニューアルと、オンラインセミナー配信環境の整備を予定しています。見積もりは3社から取得済みです。\n\nSpeaker 1 08:30\n来年度予算案を承認し、次回総会に上程することとします。次回理事会は4月15日を予定しております。\n\nSpeaker 1 09:30\n田中さん、総会資料の準備をお願いできますか。来年度予算案を含めて3月末までにお願いします。\n\nSpeaker 3 10:00\n承知しました。3月末までに準備いたします。"

        transcript_textarea = page.locator('textarea').last
        if transcript_textarea.count() > 0:
            transcript_textarea.fill(transcript_text)

        time.sleep(0.5)
        page.screenshot(path=f"{ASSETS_DIR}/03_step1_filled.png")

        # Scroll down to show transcript
        page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
        time.sleep(0.5)
        page.screenshot(path=f"{ASSETS_DIR}/03b_step1_scroll.png")

        # 5. Click Next to go to Step 2
        print("📸 Moving to Step 2...")
        page.evaluate("window.scrollTo(0, 0)")
        time.sleep(0.3)

        next_btn = page.locator('button:has-text("次へ")')
        if next_btn.count() > 0:
            next_btn.click()
            time.sleep(1)

        page.screenshot(path=f"{ASSETS_DIR}/04_step2_empty.png")

        # 6. Fill speaker names
        print("📝 Filling speaker names...")
        speaker_inputs = page.locator('input[placeholder="実名を入力..."]')
        count = speaker_inputs.count()
        names = ["山田太郎", "佐藤花子", "田中一郎"]
        for i in range(min(count, len(names))):
            speaker_inputs.nth(i).fill(names[i])
            time.sleep(0.3)

        # Click somewhere to dismiss autocomplete
        page.keyboard.press("Escape")
        time.sleep(0.5)
        page.screenshot(path=f"{ASSETS_DIR}/05_step2_filled.png")

        # 7. Generate minutes
        print("🤖 Generating minutes (this takes ~30s)...")
        gen_btn = page.locator('button:has-text("議事録を生成する")')
        if gen_btn.count() > 0:
            gen_btn.click()
            time.sleep(3)

        # Step 3 - Generating
        page.screenshot(path=f"{ASSETS_DIR}/06_step3_generating.png")

        # Wait for generation to complete
        print("⏳ Waiting for generation to complete...")
        try:
            page.wait_for_selector('text=Markdownでダウンロード', timeout=90000)
            time.sleep(1)
        except Exception as e:
            print(f"  Warning: Timeout waiting for completion: {e}")
            time.sleep(10)

        # Step 4 - Results (top)
        print("📸 Capturing Step 4 results...")
        page.evaluate("window.scrollTo(0, 0)")
        time.sleep(0.5)
        page.screenshot(path=f"{ASSETS_DIR}/07_step4_top.png")

        # Step 4 - Results (scroll to todo)
        page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
        time.sleep(0.5)
        page.screenshot(path=f"{ASSETS_DIR}/08_step4_todo.png")

        # 8. Navigate to Format Settings
        print("📸 Capturing Format Settings...")
        page.click('text=フォーマット設定')
        page.wait_for_load_state("networkidle")
        time.sleep(1)
        page.screenshot(path=f"{ASSETS_DIR}/09_format_list.png")

        # Click new template button
        new_btn = page.locator('button:has-text("新規作成"), button:has-text("テンプレート")')
        if new_btn.count() > 0:
            new_btn.first.click()
            time.sleep(1)
            page.screenshot(path=f"{ASSETS_DIR}/10_format_new.png")

        # 9. Navigate to Past Minutes
        print("📸 Capturing Past Minutes...")
        page.click('text=過去の議事録')
        page.wait_for_load_state("networkidle")
        time.sleep(1)
        page.screenshot(path=f"{ASSETS_DIR}/11_history.png")

        browser.close()
        print("✅ All screenshots captured!")


def generate_pptx():
    """Generate PPTX with embedded screenshots."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    # Constants
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)
    DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
    LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT_BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
    STEP_GREEN = RGBColor(0x05, 0x96, 0x69)
    STEP_ORANGE = RGBColor(0xEA, 0x58, 0x0C)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    def add_shape(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_rounded_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text(slide, left, top, width, height, text, size=14,
                 color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Meiryo"
        p.alignment = align
        return txBox

    def add_para(text_frame, text, size=14, color=DARK_TEXT, bold=False,
                 align=PP_ALIGN.LEFT, space_before=0):
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Meiryo"
        p.alignment = align
        p.space_before = Pt(space_before)
        return p

    def add_page_num(slide, num, total):
        add_text(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
                 f"{num} / {total}", size=9, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)

    def add_bars(slide):
        add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
        add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

    def add_header(slide, title):
        add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 title, size=26, color=PRIMARY_BLUE, bold=True)
        add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

    def add_screenshot(slide, img_path, left, top, width, height):
        """Add screenshot image, fitting within bounds while preserving aspect ratio."""
        if os.path.exists(img_path):
            from PIL import Image
            img = Image.open(img_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            target_w = width
            target_h = height
            target_aspect = target_w / target_h

            if aspect > target_aspect:
                # Image is wider - fit to width
                final_w = target_w
                final_h = int(target_w / aspect)
            else:
                # Image is taller - fit to height
                final_h = target_h
                final_w = int(target_h * aspect)

            # Center the image
            offset_x = (target_w - final_w) // 2
            offset_y = (target_h - final_h) // 2

            # Add shadow background
            shadow = add_rounded_rect(slide, left + Emu(offset_x) - Inches(0.05),
                                       top + Emu(offset_y) - Inches(0.05),
                                       Emu(final_w) + Inches(0.1),
                                       Emu(final_h) + Inches(0.1),
                                       RGBColor(0xE2, 0xE8, 0xF0))
            shadow.shadow.inherit = False

            slide.shapes.add_picture(img_path, left + Emu(offset_x), top + Emu(offset_y),
                                     Emu(final_w), Emu(final_h))
        else:
            # Placeholder if image missing
            shape = add_rounded_rect(slide, left, top, width, height, RGBColor(0xE2, 0xE8, 0xF0))
            add_text(slide, left, top + height // 2 - Inches(0.2), width, Inches(0.4),
                     f"📷 {os.path.basename(img_path)}", size=10, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    TOTAL = 11
    IMG_LEFT = Inches(7.0)
    IMG_TOP = Inches(1.3)
    IMG_W = Inches(5.8)
    IMG_H = Inches(5.5)

    # ========== Slide 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)
    add_shape(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)

    add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             "FMPJ議事録ツール", size=44, color=PRIMARY_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.6),
             "使用マニュアル", size=28, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.03), PRIMARY_BLUE)
    add_text(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
             "AI自動議事録生成システム", size=18, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    desc = add_rounded_rect(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(1.2), LIGHT_BG)
    txBox = slide.shapes.add_textbox(Inches(3.8), Inches(4.95), Inches(5.7), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PLAUDボイスレコーダーの文字起こしデータから"
    p.font.size = Pt(13); p.font.color.rgb = GRAY_TEXT; p.font.name = "Meiryo"; p.alignment = PP_ALIGN.CENTER
    add_para(tf, "Claude AIが自動で議事録・ToDoリストを作成します",
             size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER, space_before=4)

    add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
             "一般社団法人 日本音楽制作者連盟（FMPJ）", size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    # ========== Slide 2: TOC ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "📋  目次", size=28, color=PRIMARY_BLUE, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

    toc = [
        ("1", "全体フロー概要", "ツールの全体的な使い方の流れ"),
        ("2", "事前準備：PLAUDで文字起こし", "ボイスレコーダーでの録音と文字起こしデータの取得"),
        ("3", "ログイン", "ツールへのアクセスとパスワード認証"),
        ("4", "Step 1：会議情報入力", "会議名、日時、出席者、文字起こしデータの入力"),
        ("5", "Step 2：話者特定", "Speaker ラベルを実名に変換"),
        ("6", "Step 3：議事録生成", "Claude AIによる自動生成（ストリーミング）"),
        ("7", "Step 4：結果表示・ダウンロード", "議事録とToDoリストの確認・出力"),
        ("8", "フォーマット設定", "議事録フォーマットのカスタマイズ"),
        ("9", "過去の議事録", "生成済み議事録の閲覧と管理"),
    ]
    y = Inches(1.4)
    for num, title, desc in toc:
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.4), Inches(0.4))
        badge.fill.solid(); badge.fill.fore_color.rgb = PRIMARY_BLUE; badge.line.fill.background()
        tf = badge.text_frame; p = tf.paragraphs[0]
        p.text = num; p.font.size = Pt(13); p.font.color.rgb = WHITE
        p.font.bold = True; p.font.name = "Meiryo"; p.alignment = PP_ALIGN.CENTER

        add_text(slide, Inches(1.8), y - Inches(0.02), Inches(5), Inches(0.35),
                 title, size=15, color=DARK_TEXT, bold=True)
        add_text(slide, Inches(1.8), y + Inches(0.28), Inches(8), Inches(0.3),
                 desc, size=11, color=GRAY_TEXT)
        y += Inches(0.65)
    add_page_num(slide, 2, TOTAL)

    # ========== Slide 3: Flow Overview ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "1.  全体フロー概要")

    flow = [
        ("🎙️", "PLAUD録音", "会議を録音", RGBColor(0xEF, 0xF6, 0xFF)),
        ("📝", "文字起こし", "テキスト化", RGBColor(0xEF, 0xF6, 0xFF)),
        ("🔑", "ログイン", "ツールにアクセス", RGBColor(0xF0, 0xFD, 0xF4)),
        ("📋", "情報入力", "会議情報+\n文字起こし貼付", RGBColor(0xF0, 0xFD, 0xF4)),
        ("👤", "話者特定", "Speaker→実名", RGBColor(0xFF, 0xF7, 0xED)),
        ("🤖", "AI生成", "議事録を自動生成", RGBColor(0xFF, 0xF7, 0xED)),
        ("✅", "完成", "コピー/DL", RGBColor(0xEF, 0xF6, 0xFF)),
    ]
    x_start = Inches(0.5); box_w = Inches(1.55); gap = Inches(0.22); y_f = Inches(1.5)
    for i, (icon, title, desc, bg) in enumerate(flow):
        x = x_start + i * (box_w + gap)
        add_rounded_rect(slide, x, y_f, box_w, Inches(1.5), bg)
        add_text(slide, x, y_f + Inches(0.1), box_w, Inches(0.4), icon, size=22, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.1), y_f + Inches(0.5), box_w - Inches(0.2), Inches(0.3),
                 title, size=11, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
        txBox = slide.shapes.add_textbox(x + Inches(0.1), y_f + Inches(0.8), box_w - Inches(0.2), Inches(0.5))
        tf = txBox.text_frame; tf.word_wrap = True
        for line in desc.split("\n"):
            if tf.paragraphs[0].text == "": p = tf.paragraphs[0]
            else: p = tf.add_paragraph()
            p.text = line; p.font.size = Pt(9); p.font.color.rgb = GRAY_TEXT
            p.font.name = "Meiryo"; p.alignment = PP_ALIGN.CENTER
        if i < len(flow) - 1:
            add_text(slide, x + box_w + Inches(0.02), y_f + Inches(0.5), Inches(0.2), Inches(0.3),
                     "→", size=14, color=PRIMARY_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.8), Inches(3.4), Inches(5), Inches(0.4),
             "💡 ポイント", size=16, color=DARK_TEXT, bold=True)
    pts = [
        "所要時間：入力〜議事録完成まで約3〜5分",
        "データはブラウザ上で処理、サーバーには一切保存されません",
        "Markdown形式でコピー・ダウンロード可能",
        "過去の議事録はローカルストレージに自動保存",
    ]
    y_p = Inches(3.9)
    for pt in pts:
        add_text(slide, Inches(1.2), y_p, Inches(10), Inches(0.3), f"•  {pt}", size=12, color=GRAY_TEXT)
        y_p += Inches(0.38)

    # Add login screenshot on the right as a teaser
    add_screenshot(slide, f"{ASSETS_DIR}/01_login.png",
                   Inches(7.5), Inches(3.0), Inches(4.5), Inches(3.5))
    add_page_num(slide, 3, TOTAL)

    # ========== Slide 4: PLAUD ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "2.  事前準備：PLAUDで文字起こし")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "手順", size=18, color=DARK_TEXT, bold=True)

    plaud_steps = [
        ("1", "PLAUDで会議を録音", "会議開始時にPLAUDボイスレコーダーで録音を開始します。"),
        ("2", "文字起こしを実行", "録音完了後、PLAUDアプリで「文字起こし」を実行します。\n「話者分離」機能を有効にしてください。"),
        ("3", "テキストをコピー", "文字起こし結果のテキスト全体を選択してコピーします。\n「Speaker 1」等のラベル付きでコピーしてください。"),
    ]
    y_s = Inches(1.9)
    for num, title, desc in plaud_steps:
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y_s, Inches(0.35), Inches(0.35))
        badge.fill.solid(); badge.fill.fore_color.rgb = STEP_GREEN; badge.line.fill.background()
        tf = badge.text_frame; p = tf.paragraphs[0]; p.text = num
        p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
        p.font.name = "Meiryo"; p.alignment = PP_ALIGN.CENTER
        add_text(slide, Inches(1.5), y_s - Inches(0.02), Inches(5), Inches(0.3),
                 title, size=14, color=DARK_TEXT, bold=True)
        txBox = slide.shapes.add_textbox(Inches(1.5), y_s + Inches(0.3), Inches(5), Inches(0.7))
        tf2 = txBox.text_frame; tf2.word_wrap = True
        for line in desc.split("\n"):
            if tf2.paragraphs[0].text == "": p2 = tf2.paragraphs[0]
            else: p2 = tf2.add_paragraph()
            p2.text = line; p2.font.size = Pt(11); p2.font.color.rgb = GRAY_TEXT; p2.font.name = "Meiryo"
        y_s += Inches(1.2)

    tip = add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.0), RGBColor(0xFF, 0xF7, 0xED))
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(5.1), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "💡 ヒント"; p.font.size = Pt(12)
    p.font.color.rgb = STEP_ORANGE; p.font.bold = True; p.font.name = "Meiryo"
    add_para(tf, "PLAUDの「話者分離」機能でSpeaker 1, 2 等が自動付与されます。Step 2で実名に変換します。",
             size=10, color=GRAY_TEXT, space_before=4)

    # PLAUD Web screenshot
    add_screenshot(slide, f"{ASSETS_DIR}/00_plaud.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    add_page_num(slide, 4, TOTAL)

    # ========== Slide 5: Login ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "3.  ログイン")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.3),
             "アクセス方法", size=18, color=DARK_TEXT, bold=True)

    login_steps = [
        ("1", "ブラウザでアクセス", "下記URLをブラウザで開きます：\nhttps://fmpj-minutes-tool.vercel.app"),
        ("2", "パスワードを入力", "共有されたパスワードを入力します。\n👁 目のアイコンで入力内容を確認できます。"),
        ("3", "ログインをクリック", "「ログイン」ボタンをクリックすると\nダッシュボードに移動します。"),
    ]
    y_s = Inches(1.9)
    for num, title, desc in login_steps:
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y_s, Inches(0.35), Inches(0.35))
        badge.fill.solid(); badge.fill.fore_color.rgb = PRIMARY_BLUE; badge.line.fill.background()
        tf = badge.text_frame; p = tf.paragraphs[0]; p.text = num
        p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
        p.font.name = "Meiryo"; p.alignment = PP_ALIGN.CENTER
        add_text(slide, Inches(1.5), y_s - Inches(0.02), Inches(5), Inches(0.3),
                 title, size=14, color=DARK_TEXT, bold=True)
        txBox = slide.shapes.add_textbox(Inches(1.5), y_s + Inches(0.3), Inches(5), Inches(0.6))
        tf2 = txBox.text_frame; tf2.word_wrap = True
        for line in desc.split("\n"):
            if tf2.paragraphs[0].text == "": p2 = tf2.paragraphs[0]
            else: p2 = tf2.add_paragraph()
            p2.text = line; p2.font.size = Pt(11); p2.font.color.rgb = GRAY_TEXT; p2.font.name = "Meiryo"
        y_s += Inches(1.05)

    note = add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.7), LIGHT_BG)
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(5.1), Inches(0.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "ℹ️ セッションは24時間有効です。期限切れ後は再ログインが必要です。"
    p.font.size = Pt(11); p.font.color.rgb = GRAY_TEXT; p.font.name = "Meiryo"

    add_screenshot(slide, f"{ASSETS_DIR}/01_login.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)
    add_page_num(slide, 5, TOTAL)

    # ========== Slide 6: Step 1 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "4.  Step 1：会議情報入力")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "入力項目", size=18, color=DARK_TEXT, bold=True)

    fields = [
        ("会議名", "例：第10回理事会"),
        ("会議種別", "理事会 / 委員会 / 総会 / 部会 / WG / その他"),
        ("開催日", "会議の開催日を入力"),
        ("開催場所", "例：連盟会議室 / オンライン"),
        ("出席者", "出席者名を改行区切りで入力"),
        ("フォーマット", "事前設定したテンプレートを選択（任意）"),
        ("文字起こし", "PLAUDからコピーしたテキストを貼り付け"),
    ]
    y_s = Inches(1.85)
    for field, desc in fields:
        add_text(slide, Inches(1.0), y_s, Inches(1.6), Inches(0.25),
                 f"▸  {field}", size=11, color=DARK_TEXT, bold=True)
        add_text(slide, Inches(2.7), y_s, Inches(3.7), Inches(0.25), desc, size=10, color=GRAY_TEXT)
        y_s += Inches(0.36)

    tip = add_rounded_rect(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(1.6), RGBColor(0xF0, 0xFD, 0xF4))
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(5.1), Inches(1.4))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "💡 ヒント"; p.font.size = Pt(12)
    p.font.color.rgb = STEP_GREEN; p.font.bold = True; p.font.name = "Meiryo"
    add_para(tf, "• 出席者名はStep 2の話者特定候補に使われます", size=10, color=GRAY_TEXT, space_before=6)
    add_para(tf, "• テンプレートで出力形式をカスタマイズ可能", size=10, color=GRAY_TEXT, space_before=4)
    add_para(tf, "• 入力完了後「次へ」をクリック", size=10, color=GRAY_TEXT, space_before=4)

    add_screenshot(slide, f"{ASSETS_DIR}/03_step1_filled.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)
    add_page_num(slide, 6, TOTAL)

    # ========== Slide 7: Step 2 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "5.  Step 2：話者特定")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "操作方法", size=18, color=DARK_TEXT, bold=True)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.6))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "文字起こしデータから「Speaker 1」「Speaker 2」等が自動検出されます。"
    p.font.size = Pt(12); p.font.color.rgb = GRAY_TEXT; p.font.name = "Meiryo"

    add_text(slide, Inches(0.8), Inches(2.6), Inches(3), Inches(0.3),
             "割り当て方法", size=14, color=DARK_TEXT, bold=True)

    methods = [
        ("方法1：出席者チップをクリック", "画面上部の名前チップをクリックで自動入力"),
        ("方法2：直接入力", "入力欄に実名を直接入力（サジェスト付き）"),
    ]
    y_s = Inches(3.0)
    for title, desc in methods:
        add_text(slide, Inches(1.0), y_s, Inches(5), Inches(0.3), title, size=12, color=DARK_TEXT, bold=True)
        add_text(slide, Inches(1.0), y_s + Inches(0.28), Inches(5), Inches(0.3), desc, size=10, color=GRAY_TEXT)
        y_s += Inches(0.7)

    add_text(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.5),
             "💡 各Speakerの横の「(◯回発言)」で発言頻度を確認し、正しい人物を特定してください。",
             size=10, color=GRAY_TEXT)

    note = add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.8), RGBColor(0xFF, 0xF7, 0xED))
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(5.1), Inches(0.6))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "⚠️ すべてのSpeakerに実名を割り当ててから「議事録を生成する」をクリック"
    p.font.size = Pt(11); p.font.color.rgb = STEP_ORANGE; p.font.name = "Meiryo"

    add_screenshot(slide, f"{ASSETS_DIR}/05_step2_filled.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)
    add_page_num(slide, 7, TOTAL)

    # ========== Slide 8: Step 3 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "6.  Step 3：議事録生成")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "生成プロセス", size=18, color=DARK_TEXT, bold=True)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.5), Inches(1.2))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Claude AI（Anthropic社）が文字起こしデータを分析し、議事録とToDoリストを自動生成します。"
    p.font.size = Pt(12); p.font.color.rgb = GRAY_TEXT; p.font.name = "Meiryo"
    add_para(tf, "", size=4)
    add_para(tf, "生成内容：", size=13, color=DARK_TEXT, bold=True, space_before=8)
    add_para(tf, "  ✅  議事録（ヘッダー + 議事内容 + 決定事項）", size=11, color=GRAY_TEXT, space_before=4)
    add_para(tf, "  ✅  ToDoリスト（担当者・タスク・期限）", size=11, color=GRAY_TEXT, space_before=2)

    feat = add_rounded_rect(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.2), LIGHT_BG)
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(5.1), Inches(2.0))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "🔧 生成中の操作"; p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT; p.font.bold = True; p.font.name = "Meiryo"
    add_para(tf, "• ストリーミング表示でリアルタイムに確認可能", size=11, color=GRAY_TEXT, space_before=8)
    add_para(tf, "• 所要時間：約30秒〜1分（1時間の会議）", size=11, color=GRAY_TEXT, space_before=4)
    add_para(tf, "• キャンセルボタンで途中中断可能", size=11, color=GRAY_TEXT, space_before=4)
    add_para(tf, "• 完了後、自動でStep 4に移動", size=11, color=GRAY_TEXT, space_before=4)

    add_screenshot(slide, f"{ASSETS_DIR}/06_step3_generating.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)
    add_page_num(slide, 8, TOTAL)

    # ========== Slide 9: Step 4 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "7.  Step 4：結果表示・ダウンロード")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "出力内容と操作", size=18, color=DARK_TEXT, bold=True)

    actions = [
        ("📋", "コピー", "各セクションの「コピー」ボタンで\nクリップボードにコピー"),
        ("⬇️", "ダウンロード", "「Markdownでダウンロード」で\n.mdファイルとして保存"),
        ("🔄", "やり直し", "「最初からやり直す」で\nStep 1に戻って新規作成"),
    ]
    y_s = Inches(1.9)
    for icon, title, desc in actions:
        box = add_rounded_rect(slide, Inches(0.8), y_s, Inches(5.5), Inches(1.1), LIGHT_BG)
        add_text(slide, Inches(1.0), y_s + Inches(0.1), Inches(0.4), Inches(0.3), icon, size=16, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.5), y_s + Inches(0.1), Inches(1.5), Inches(0.3), title, size=13, color=DARK_TEXT, bold=True)
        txBox = slide.shapes.add_textbox(Inches(1.5), y_s + Inches(0.42), Inches(4.5), Inches(0.6))
        tf = txBox.text_frame; tf.word_wrap = True
        for line in desc.split("\n"):
            if tf.paragraphs[0].text == "": p = tf.paragraphs[0]
            else: p = tf.add_paragraph()
            p.text = line; p.font.size = Pt(10); p.font.color.rgb = GRAY_TEXT; p.font.name = "Meiryo"
        y_s += Inches(1.25)

    note = add_rounded_rect(slide, Inches(0.8), Inches(5.7), Inches(5.5), Inches(0.6), RGBColor(0xF0, 0xFD, 0xF4))
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.78), Inches(5.1), Inches(0.4))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "✅ 生成結果は「過去の議事録」に自動保存されます"
    p.font.size = Pt(11); p.font.color.rgb = STEP_GREEN; p.font.name = "Meiryo"

    add_screenshot(slide, f"{ASSETS_DIR}/07_step4_top.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)
    add_page_num(slide, 9, TOTAL)

    # ========== Slide 10: Format Settings ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "8.  フォーマット設定")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "テンプレート管理", size=18, color=DARK_TEXT, bold=True)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(5.5), Inches(0.6))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ナビゲーションの「フォーマット設定」から、議事録フォーマットをカスタマイズできます。"
    p.font.size = Pt(12); p.font.color.rgb = GRAY_TEXT; p.font.name = "Meiryo"

    settings = [
        ("テンプレート名", "例：理事会用フォーマット"),
        ("会議種別", "対象の会議種別を選択"),
        ("フォーマット指示", "追加の書式ルール"),
        ("サンプル出力", "過去の議事録（Few-shot学習用）"),
    ]
    y_s = Inches(2.7)
    for field, desc in settings:
        add_text(slide, Inches(1.0), y_s, Inches(2.0), Inches(0.25),
                 f"▸  {field}", size=11, color=DARK_TEXT, bold=True)
        add_text(slide, Inches(3.1), y_s, Inches(3.2), Inches(0.25), desc, size=10, color=GRAY_TEXT)
        y_s += Inches(0.35)

    tip = add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.0), RGBColor(0xF0, 0xFD, 0xF4))
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(5.1), Inches(1.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "💡 サンプル出力の活用"; p.font.size = Pt(12)
    p.font.color.rgb = STEP_GREEN; p.font.bold = True; p.font.name = "Meiryo"
    add_para(tf, "過去の議事録をサンプルとして設定すると、AIがそのスタイルを学習し同様の形式で生成します。",
             size=10, color=GRAY_TEXT, space_before=6)
    add_para(tf, "初回は手動作成の議事録を設定→2回目以降はAI生成の議事録をそのまま利用できます。",
             size=10, color=GRAY_TEXT, space_before=4)

    add_screenshot(slide, f"{ASSETS_DIR}/10_format_new.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)
    add_page_num(slide, 10, TOTAL)

    # ========== Slide 11: Past Minutes ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bars(slide); add_header(slide, "9.  過去の議事録")

    add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "履歴管理", size=18, color=DARK_TEXT, bold=True)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(5.5), Inches(0.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "「過去の議事録」で生成済み議事録の閲覧・検索・出力が可能です。"
    p.font.size = Pt(12); p.font.color.rgb = GRAY_TEXT; p.font.name = "Meiryo"

    features = [
        ("🏷️", "会議種別フィルター", "カテゴリで絞り込み"),
        ("🔍", "テキスト検索", "会議名・内容・出席者で検索"),
        ("📖", "内容表示", "カードクリックで全文展開"),
        ("📋", "コピー", "クリップボードにコピー"),
        ("⬇️", "ダウンロード", "Markdownファイルで保存"),
        ("🗑️", "削除", "不要な議事録を削除"),
    ]
    y_s = Inches(2.5)
    for icon, title, desc in features:
        add_text(slide, Inches(1.0), y_s, Inches(0.3), Inches(0.25), icon, size=12)
        add_text(slide, Inches(1.35), y_s, Inches(1.8), Inches(0.25), title, size=11, color=DARK_TEXT, bold=True)
        add_text(slide, Inches(3.3), y_s, Inches(3.0), Inches(0.25), desc, size=10, color=GRAY_TEXT)
        y_s += Inches(0.36)

    note = add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(1.2), RGBColor(0xFF, 0xF7, 0xED))
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(5.1), Inches(1.0))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "⚠️ データ保存について"; p.font.size = Pt(12)
    p.font.color.rgb = STEP_ORANGE; p.font.bold = True; p.font.name = "Meiryo"
    add_para(tf, "データはブラウザのローカルストレージに保存されます。ブラウザのキャッシュ削除で履歴も消えますのでご注意ください。重要な議事録は必ずダウンロードして保管してください。",
             size=10, color=GRAY_TEXT, space_before=4)

    add_screenshot(slide, f"{ASSETS_DIR}/11_history.png", IMG_LEFT, IMG_TOP, IMG_W, IMG_H)
    add_page_num(slide, 11, TOTAL)

    # Save
    output = "/Users/akirahasama/Desktop/dev/fmpj-minutes-tool/FMPJ議事録ツール_使用マニュアル.pptx"
    prs.save(output)
    print(f"✅ PPTX saved: {output}")


if __name__ == "__main__":
    print("=" * 50)
    print("FMPJ議事録ツール マニュアル生成")
    print("=" * 50)
    print()
    print("Phase 1: Capturing screenshots...")
    capture_screenshots()
    print()
    print("Phase 2: Generating PPTX...")
    generate_pptx()
    print()
    print("🎉 Done!")
