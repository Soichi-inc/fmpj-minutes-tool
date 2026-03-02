#!/usr/bin/env python3
"""FMPJ議事録ツール 使用マニュアル PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
PLACEHOLDER_BG = RGBColor(0xE2, 0xE8, 0xF0)
STEP_GREEN = RGBColor(0x05, 0x96, 0x69)
STEP_ORANGE = RGBColor(0xEA, 0x58, 0x0C)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color=WHITE):
    """Add background color to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, border_color=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=DARK_TEXT, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name="Meiryo"):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_image_placeholder(slide, left, top, width, height, label="スクリーンショット"):
    """Add an image placeholder with label."""
    shape = add_shape(slide, left, top, width, height, PLACEHOLDER_BG, PRIMARY_BLUE)
    # Add label text on the placeholder
    txBox = slide.shapes.add_textbox(left, top + height // 2 - Inches(0.3), width, Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"📷 {label}"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Meiryo"
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_step_badge(slide, left, top, step_num, label, color=PRIMARY_BLUE):
    """Add a step number badge."""
    # Circle with number
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(step_num)
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Meiryo"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)

    # Label
    add_text_box(slide, left + Inches(0.65), top + Inches(0.05), Inches(3), Inches(0.4),
                 label, font_size=16, color=DARK_TEXT, bold=True)


def add_page_number(slide, num, total):
    """Add page number to bottom right."""
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
                 f"{num} / {total}", font_size=9, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)


# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, WHITE)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)

# Title area
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             "FMPJ議事録ツール", font_size=44, color=PRIMARY_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.6),
             "使用マニュアル", font_size=28, color=DARK_TEXT, bold=False,
             alignment=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.03), PRIMARY_BLUE)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
             "AI自動議事録生成システム", font_size=18, color=GRAY_TEXT,
             alignment=PP_ALIGN.CENTER)

# Description box
desc_box = add_rounded_rect(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(1.2), LIGHT_BG)
txBox = slide.shapes.add_textbox(Inches(3.8), Inches(4.95), Inches(5.7), Inches(0.9))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PLAUDボイスレコーダーの文字起こしデータから"
p.font.size = Pt(13)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"
p.alignment = PP_ALIGN.CENTER
add_paragraph(tf, "Claude AIが自動で議事録・ToDoリストを作成します",
              font_size=13, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER, space_before=4)

# Footer
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
             "一般社団法人 日本音楽制作者連盟（FMPJ）", font_size=12, color=GRAY_TEXT,
             alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)


# ============================================================
# Slide 2: Table of Contents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "📋  目次", font_size=28, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

toc_items = [
    ("1", "全体フロー概要", "ツールの全体的な使い方の流れ"),
    ("2", "事前準備：PLAUDで文字起こし", "ボイスレコーダーでの録音と文字起こしデータの取得"),
    ("3", "ログイン", "ツールへのアクセスとパスワード認証"),
    ("4", "Step 1：会議情報入力", "会議名、日時、出席者、文字起こしデータの入力"),
    ("5", "Step 2：話者特定", "Speaker ラベルを実名に変換"),
    ("6", "Step 3：議事録生成", "Claude AIによる自動生成（ストリーミング）"),
    ("7", "Step 4：結果表示・ダウンロード", "議事録とToDoリストの確認・出力"),
    ("8", "フォーマット設定", "議事録フォーマットのカスタマイズ"),
    ("9", "過去の議事録", "生成済み議事録の閲覧と管理"),
]

y_pos = Inches(1.4)
for num, title, desc in toc_items:
    # Number circle
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y_pos, Inches(0.4), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PRIMARY_BLUE
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Meiryo"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.8), y_pos - Inches(0.02), Inches(5), Inches(0.35),
                 title, font_size=15, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(1.8), y_pos + Inches(0.28), Inches(8), Inches(0.3),
                 desc, font_size=11, color=GRAY_TEXT)

    y_pos += Inches(0.65)

add_page_number(slide, 2, 11)


# ============================================================
# Slide 3: Overall Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "1.  全体フロー概要", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Flow boxes
flow_items = [
    ("🎙️", "PLAUD録音", "会議をPLAUDで\n録音", RGBColor(0xEF, 0xF6, 0xFF)),
    ("📝", "文字起こし", "PLAUDアプリで\nテキスト化", RGBColor(0xEF, 0xF6, 0xFF)),
    ("🔑", "ログイン", "パスワードで\nツールにアクセス", RGBColor(0xF0, 0xFD, 0xF4)),
    ("📋", "情報入力", "会議情報と\n文字起こし貼付", RGBColor(0xF0, 0xFD, 0xF4)),
    ("👤", "話者特定", "Speaker→\n実名に変換", RGBColor(0xFF, 0xF7, 0xED)),
    ("🤖", "AI生成", "Claude AIが\n議事録を自動生成", RGBColor(0xFF, 0xF7, 0xED)),
    ("✅", "完成", "議事録・ToDo\nコピー/DL", RGBColor(0xEF, 0xF6, 0xFF)),
]

x_start = Inches(0.5)
box_w = Inches(1.55)
box_h = Inches(1.8)
gap = Inches(0.22)
y_flow = Inches(1.6)

for i, (icon, title, desc, bg_color) in enumerate(flow_items):
    x = x_start + i * (box_w + gap)

    box = add_rounded_rect(slide, x, y_flow, box_w, box_h, bg_color)

    # Icon
    add_text_box(slide, x, y_flow + Inches(0.15), box_w, Inches(0.4),
                 icon, font_size=24, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Inches(0.1), y_flow + Inches(0.6), box_w - Inches(0.2), Inches(0.3),
                 title, font_size=12, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    # Desc
    txBox = slide.shapes.add_textbox(x + Inches(0.1), y_flow + Inches(0.95),
                                      box_w - Inches(0.2), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        if tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY_TEXT
        p.font.name = "Meiryo"
        p.alignment = PP_ALIGN.CENTER

    # Arrow (except last)
    if i < len(flow_items) - 1:
        arrow_x = x + box_w + Inches(0.02)
        add_text_box(slide, arrow_x, y_flow + Inches(0.7), Inches(0.2), Inches(0.3),
                     "→", font_size=16, color=PRIMARY_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Key points
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.4),
             "💡 ポイント", font_size=16, color=DARK_TEXT, bold=True)

points = [
    "所要時間：入力〜議事録完成まで約3〜5分（1時間の会議の場合）",
    "データはブラウザ上で処理され、サーバーには一切保存されません",
    "生成された議事録はMarkdown形式でコピー・ダウンロード可能",
    "過去の議事録はブラウザのローカルストレージに自動保存されます",
]

y_pt = Inches(4.3)
for pt in points:
    add_text_box(slide, Inches(1.2), y_pt, Inches(10), Inches(0.3),
                 f"•  {pt}", font_size=12, color=GRAY_TEXT)
    y_pt += Inches(0.38)

add_page_number(slide, 3, 11)


# ============================================================
# Slide 4: PLAUD Preparation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "2.  事前準備：PLAUDで文字起こし", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left: Instructions
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
             "手順", font_size=18, color=DARK_TEXT, bold=True)

steps = [
    ("1", "PLAUDで会議を録音", "会議開始時にPLAUDボイスレコーダーで録音を開始します。"),
    ("2", "文字起こしを実行", "録音完了後、PLAUDアプリで「文字起こし」を実行します。"),
    ("3", "テキストをコピー", "文字起こし結果のテキスト全体を選択してコピーします。\n「Speaker 1」「Speaker 2」等のラベルが付いた状態でコピーしてください。"),
]

y_s = Inches(1.9)
for num, title, desc in steps:
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y_s, Inches(0.35), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = STEP_GREEN
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Meiryo"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.5), y_s - Inches(0.02), Inches(5), Inches(0.3),
                 title, font_size=14, color=DARK_TEXT, bold=True)

    txBox = slide.shapes.add_textbox(Inches(1.5), y_s + Inches(0.3), Inches(5), Inches(0.8))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    for line in desc.split("\n"):
        if tf2.paragraphs[0].text == "":
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY_TEXT
        p2.font.name = "Meiryo"

    y_s += Inches(1.2)

# Tip box
tip = add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(5.8), Inches(1.2), RGBColor(0xFF, 0xF7, 0xED))
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(5.4), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "💡 ヒント"
p.font.size = Pt(12)
p.font.color.rgb = STEP_ORANGE
p.font.bold = True
p.font.name = "Meiryo"
add_paragraph(tf, "PLAUDアプリの「話者分離」機能を有効にすると、Speaker 1, Speaker 2 のように自動で話者が分離されます。この情報をStep 2で実名に変換します。",
              font_size=10, color=GRAY_TEXT, space_before=4)

# Right: Screenshot placeholders
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                      "PLAUDアプリ画面のスクリーンショットを挿入")

add_page_number(slide, 4, 11)


# ============================================================
# Slide 5: Login
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "3.  ログイン", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left content
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
             "アクセス方法", font_size=18, color=DARK_TEXT, bold=True)

login_steps = [
    ("1", "ブラウザでアクセス", "下記URLをブラウザで開きます：\nhttps://fmpj-minutes-tool.vercel.app"),
    ("2", "パスワードを入力", "共有されたパスワードを入力します。\n👁 目のアイコンで入力内容を確認できます。"),
    ("3", "ログインをクリック", "「ログイン」ボタンをクリックするとダッシュボードに移動します。"),
]

y_s = Inches(1.9)
for num, title, desc in login_steps:
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y_s, Inches(0.35), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PRIMARY_BLUE
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Meiryo"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.5), y_s - Inches(0.02), Inches(5), Inches(0.3),
                 title, font_size=14, color=DARK_TEXT, bold=True)

    txBox = slide.shapes.add_textbox(Inches(1.5), y_s + Inches(0.3), Inches(5), Inches(0.7))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    for line in desc.split("\n"):
        if tf2.paragraphs[0].text == "":
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY_TEXT
        p2.font.name = "Meiryo"

    y_s += Inches(1.1)

# Note box
note = add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(5.8), Inches(0.8), LIGHT_BG)
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(5.4), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ℹ️ セッションは24時間有効です。期限切れ後は再ログインが必要です。"
p.font.size = Pt(11)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"

# Right: Screenshot
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.0),
                      "ログイン画面のスクリーンショットを挿入")

add_page_number(slide, 5, 11)


# ============================================================
# Slide 6: Step 1 - Meeting Info Input
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "4.  Step 1：会議情報入力", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left content
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "入力項目", font_size=18, color=DARK_TEXT, bold=True)

fields = [
    ("会議名", "例：第10回理事会"),
    ("会議種別", "理事会 / 委員会 / 総会 / 部会 / ワーキンググループ / その他"),
    ("開催日", "会議の開催日を入力"),
    ("開催場所", "例：連盟会議室 / オンライン"),
    ("出席者", "出席者名を改行区切りで入力（話者特定の候補になります）"),
    ("フォーマット", "事前設定したテンプレートを選択（任意）"),
    ("文字起こしデータ", "PLAUDからコピーしたテキストを貼り付け"),
]

y_s = Inches(1.85)
for field, desc in fields:
    add_text_box(slide, Inches(1.0), y_s, Inches(1.8), Inches(0.25),
                 f"▸  {field}", font_size=11, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(2.9), y_s, Inches(3.5), Inches(0.25),
                 desc, font_size=10, color=GRAY_TEXT)
    y_s += Inches(0.38)

# Tip
tip = add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(5.8), Inches(1.5), RGBColor(0xF0, 0xFD, 0xF4))
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(5.4), Inches(1.3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "💡 ヒント"
p.font.size = Pt(12)
p.font.color.rgb = STEP_GREEN
p.font.bold = True
p.font.name = "Meiryo"
add_paragraph(tf, "• 出席者名は正確に入力してください。Step 2 での話者特定に使われます。",
              font_size=10, color=GRAY_TEXT, space_before=6)
add_paragraph(tf, "• フォーマットテンプレートを事前に設定しておくと、会議種別に応じた出力形式を指定できます。",
              font_size=10, color=GRAY_TEXT, space_before=4)
add_paragraph(tf, "• 入力が完了したら「次へ」ボタンをクリックします。",
              font_size=10, color=GRAY_TEXT, space_before=4)

# Right: Screenshot
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                      "Step 1 会議情報入力画面を挿入")

add_page_number(slide, 6, 11)


# ============================================================
# Slide 7: Step 2 - Speaker Identification
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "5.  Step 2：話者特定", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left content
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "操作方法", font_size=18, color=DARK_TEXT, bold=True)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "文字起こしデータから「Speaker 1」「Speaker 2」などの話者ラベルが自動検出されます。"
p.font.size = Pt(12)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"

add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "話者の割り当て方法：", font_size=13, color=DARK_TEXT, bold=True, space_before=8)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "方法1：出席者チップをクリック", font_size=12, color=DARK_TEXT, bold=True, space_before=6)
add_paragraph(tf, "  画面上部の出席者名チップをクリックすると、対応する入力欄に名前が自動入力されます。",
              font_size=10, color=GRAY_TEXT, space_before=2)
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "方法2：直接入力", font_size=12, color=DARK_TEXT, bold=True, space_before=6)
add_paragraph(tf, "  各Speakerの入力欄に実名を直接入力します。入力時にサジェストが表示されます。",
              font_size=10, color=GRAY_TEXT, space_before=2)
add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "各Speakerの横に表示される「(◯回発言)」は、その話者の発言回数です。発言内容と照らし合わせて正しい人物を特定してください。",
              font_size=10, color=GRAY_TEXT, space_before=8)

# Note
note = add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(5.8), Inches(0.9), RGBColor(0xFF, 0xF7, 0xED))
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(5.4), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "⚠️ 注意"
p.font.size = Pt(12)
p.font.color.rgb = STEP_ORANGE
p.font.bold = True
p.font.name = "Meiryo"
add_paragraph(tf, "すべてのSpeakerに実名を割り当ててから「議事録を生成する」ボタンをクリックしてください。",
              font_size=10, color=GRAY_TEXT, space_before=4)

# Right: Screenshot
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                      "Step 2 話者特定画面を挿入")

add_page_number(slide, 7, 11)


# ============================================================
# Slide 8: Step 3 - Generation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "6.  Step 3：議事録生成", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left content
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "生成プロセス", font_size=18, color=DARK_TEXT, bold=True)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "「議事録を生成する」ボタンをクリックすると、Claude AI（Anthropic社）が文字起こしデータを分析し、議事録とToDoリストを自動生成します。"
p.font.size = Pt(12)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"

add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "生成内容：", font_size=13, color=DARK_TEXT, bold=True, space_before=8)
add_paragraph(tf, "  ✅  議事録（会議ヘッダー + 議事内容 + 決定事項）", font_size=11, color=GRAY_TEXT, space_before=4)
add_paragraph(tf, "  ✅  ToDoリスト（担当者・タスク・期限）", font_size=11, color=GRAY_TEXT, space_before=2)

add_paragraph(tf, "", font_size=6)
add_paragraph(tf, "生成はストリーミング形式で行われ、リアルタイムで文章が表示されます。1時間の会議で約30秒〜1分程度です。",
              font_size=11, color=GRAY_TEXT, space_before=6)

# Features box
feat = add_rounded_rect(slide, Inches(0.8), Inches(4.5), Inches(5.8), Inches(1.8), LIGHT_BG)
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(5.4), Inches(1.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🔧 生成中の操作"
p.font.size = Pt(13)
p.font.color.rgb = DARK_TEXT
p.font.bold = True
p.font.name = "Meiryo"
add_paragraph(tf, "", font_size=4)
add_paragraph(tf, "• キャンセルボタン：生成を途中で中断できます", font_size=11, color=GRAY_TEXT, space_before=4)
add_paragraph(tf, "• 生成完了後、自動的にStep 4（結果表示）に移動します", font_size=11, color=GRAY_TEXT, space_before=4)
add_paragraph(tf, "• 生成結果はブラウザの履歴に自動保存されます", font_size=11, color=GRAY_TEXT, space_before=4)

# Right: Screenshot
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                      "Step 3 生成中画面を挿入")

add_page_number(slide, 8, 11)


# ============================================================
# Slide 9: Step 4 - Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "7.  Step 4：結果表示・ダウンロード", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left content
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "出力内容と操作", font_size=18, color=DARK_TEXT, bold=True)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "生成完了後、議事録とToDoリストが別セクションで表示されます。"
p.font.size = Pt(12)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"

# Action boxes
actions = [
    ("📋", "コピー", "各セクションの「コピー」ボタンで\nクリップボードにコピーします。\nメールやドキュメントに貼り付けできます。"),
    ("⬇️", "ダウンロード", "「Markdownでダウンロード」ボタンで\n議事録全体を .md ファイルとして\nダウンロードします。"),
    ("🔄", "やり直し", "「最初からやり直す」ボタンで\nStep 1に戻り、新しい議事録を\n作成できます。"),
]

y_s = Inches(2.8)
for icon, title, desc in actions:
    box = add_rounded_rect(slide, Inches(0.8), y_s, Inches(5.8), Inches(1.2), LIGHT_BG)

    add_text_box(slide, Inches(1.0), y_s + Inches(0.1), Inches(0.5), Inches(0.3),
                 icon, font_size=18, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), y_s + Inches(0.1), Inches(1.5), Inches(0.3),
                 title, font_size=13, color=DARK_TEXT, bold=True)

    txBox = slide.shapes.add_textbox(Inches(1.5), y_s + Inches(0.45), Inches(4.8), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        if tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY_TEXT
        p.font.name = "Meiryo"

    y_s += Inches(1.35)

# Right: Screenshot
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                      "Step 4 結果表示画面を挿入")

add_page_number(slide, 9, 11)


# ============================================================
# Slide 10: Format Settings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "8.  フォーマット設定", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left content
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "テンプレート管理", font_size=18, color=DARK_TEXT, bold=True)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ナビゲーションバーの「フォーマット設定」からアクセスします。会議種別ごとに議事録のフォーマットをカスタマイズできます。"
p.font.size = Pt(12)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"

# Settings fields
add_text_box(slide, Inches(0.8), Inches(2.7), Inches(3), Inches(0.3),
             "設定項目", font_size=14, color=DARK_TEXT, bold=True)

settings = [
    ("テンプレート名", "わかりやすい名前（例：理事会用フォーマット）"),
    ("会議種別", "対象の会議種別を選択"),
    ("説明", "テンプレートの説明（任意）"),
    ("フォーマット指示", "議事録に反映したい追加の書式ルール"),
    ("サンプル出力", "参考となる過去の議事録を貼付（Few-shot学習）"),
]

y_s = Inches(3.1)
for field, desc in settings:
    add_text_box(slide, Inches(1.0), y_s, Inches(2.0), Inches(0.25),
                 f"▸  {field}", font_size=11, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(3.1), y_s, Inches(3.3), Inches(0.25),
                 desc, font_size=10, color=GRAY_TEXT)
    y_s += Inches(0.35)

# Tip
tip = add_rounded_rect(slide, Inches(0.8), Inches(5.1), Inches(5.8), Inches(1.4), RGBColor(0xF0, 0xFD, 0xF4))
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(5.4), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "💡 サンプル出力の活用"
p.font.size = Pt(12)
p.font.color.rgb = STEP_GREEN
p.font.bold = True
p.font.name = "Meiryo"
add_paragraph(tf, "過去の議事録をサンプルとして設定すると、AIがそのスタイルを学習し、同様の形式で新しい議事録を生成します。初回は手動で作成した議事録を設定し、2回目以降はAI生成の議事録をそのまま利用できます。",
              font_size=10, color=GRAY_TEXT, space_before=4)

# Right: Screenshot
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                      "フォーマット設定画面を挿入")

add_page_number(slide, 10, 11)


# ============================================================
# Slide 11: Past Minutes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)
add_shape(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), PRIMARY_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "9.  過去の議事録", font_size=26, color=PRIMARY_BLUE, bold=True)

add_shape(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.02), ACCENT_BLUE_LIGHT)

# Left content
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.3),
             "履歴管理", font_size=18, color=DARK_TEXT, bold=True)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.8), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ナビゲーションバーの「過去の議事録」からアクセスします。生成した議事録の一覧を確認・検索・出力できます。"
p.font.size = Pt(12)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"

# Features
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(3), Inches(0.3),
             "機能一覧", font_size=14, color=DARK_TEXT, bold=True)

features = [
    ("🏷️", "会議種別フィルター", "理事会、委員会などのカテゴリで絞り込み"),
    ("🔍", "テキスト検索", "会議名・内容・出席者名で検索"),
    ("📖", "内容表示", "カードをクリックして議事録全文を展開"),
    ("📋", "コピー", "過去の議事録をクリップボードにコピー"),
    ("⬇️", "ダウンロード", "Markdownファイルとしてダウンロード"),
    ("🗑️", "削除", "不要な議事録を削除"),
]

y_s = Inches(3.2)
for icon, title, desc in features:
    add_text_box(slide, Inches(1.0), y_s, Inches(0.3), Inches(0.25), icon, font_size=12)
    add_text_box(slide, Inches(1.35), y_s, Inches(1.8), Inches(0.25),
                 title, font_size=11, color=DARK_TEXT, bold=True)
    add_text_box(slide, Inches(3.3), y_s, Inches(3.2), Inches(0.25),
                 desc, font_size=10, color=GRAY_TEXT)
    y_s += Inches(0.38)

# Note
note = add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(5.8), Inches(0.8), RGBColor(0xFF, 0xF7, 0xED))
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.7), Inches(5.4), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "⚠️ データはブラウザのローカルストレージに保存されています。ブラウザのキャッシュを削除すると履歴も消えますのでご注意ください。"
p.font.size = Pt(10)
p.font.color.rgb = GRAY_TEXT
p.font.name = "Meiryo"

# Right: Screenshot
add_image_placeholder(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                      "過去の議事録画面を挿入")

add_page_number(slide, 11, 11)


# ============================================================
# Save
# ============================================================
output_path = "/Users/akirahasama/Desktop/dev/fmpj-minutes-tool/FMPJ議事録ツール_使用マニュアル.pptx"
prs.save(output_path)
print(f"✅ PPTX generated: {output_path}")
